#!/usr/bin/env python3
"""
Script to analyze PPTX presentation for layout issues
"""

from pptx import Presentation
from pptx.util import Inches, Pt
import os

def analyze_presentation(pptx_path):
    """Analyze a PPTX file for layout issues"""

    if not os.path.exists(pptx_path):
        print(f"Error: File '{pptx_path}' does not exist")
        return

    print(f"Analyzing presentation: {pptx_path}")
    print("=" * 80)

    # Slide dimensions (16:9 layout)
    SLIDE_WIDTH = 10.0  # inches
    SLIDE_HEIGHT = 5.625  # inches

    prs = Presentation(pptx_path)

    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"\nSlide {slide_num}:")
        print("-" * 40)

        # Get slide dimensions
        slide_width = prs.slide_width / 914400  # Convert EMU to inches
        slide_height = prs.slide_height / 914400  # Convert EMU to inches

        print(f"Slide dimensions: {slide_width:.3f} x {slide_height:.3f} inches")

        # Check all shapes on the slide
        shapes_issues = []

        for shape_idx, shape in enumerate(shape for shape in slide.shapes if shape.has_text_frame or shape.has_table):
            shape_name = f"Shape {shape_idx + 1}"

            # Get shape position and size
            left = shape.left / 914400 if shape.left else 0
            top = shape.top / 914400 if shape.top else 0
            width = shape.width / 914400 if shape.width else 0
            height = shape.height / 914400 if shape.height else 0

            # Check 1: Boundaries
            if left + width > SLIDE_WIDTH:
                shapes_issues.append(f"{shape_name}: Extends beyond right boundary (x+w={left+width:.3f} > {SLIDE_WIDTH})")

            if top + height > SLIDE_HEIGHT:
                shapes_issues.append(f"{shape_name}: Extends beyond bottom boundary (y+h={top+height:.3f} > {SLIDE_HEIGHT})")

            # Check text content
            if shape.has_text_frame:
                text_frame = shape.text_frame

                # Check 3: Font size
                for para_idx, paragraph in enumerate(text_frame.paragraphs):
                    for run_idx, run in enumerate(paragraph.runs):
                        if run.font.size:
                            font_size = Pt(run.font.size)
                            if font_size < Pt(8):
                                shapes_issues.append(f"{shape_name} - Paragraph {para_idx+1}, Run {run_idx+1}: Font size too small ({font_size:.1f}pt < 8pt)")

            # Check 4: Tables
            elif shape.has_table:
                table = shape.table
                total_col_width = 0

                # Calculate total column width
                for col_idx, column in enumerate(table.columns):
                    col_width = column.width / 914400
                    total_col_width += col_width

                # Check if column widths sum to table width (with tolerance)
                if abs(total_col_width - width) > 0.1:  # 0.1 inch tolerance
                    shapes_issues.append(f"{shape_name}: Table column widths don't match table width (sum={total_col_width:.3f}, table width={width:.3f})")

        # Check for overlapping shapes
        shapes_list = [(shape.left/914400, shape.top/914400, shape.width/914400, shape.height/914400,
                       f"Shape {i+1}")
                      for i, shape in enumerate(slide.shapes)
                      if shape.has_text_frame or shape.has_table]

        overlapping_shapes = []
        for i in range(len(shapes_list)):
            for j in range(i + 1, len(shapes_list)):
                x1, y1, w1, h1, name1 = shapes_list[i]
                x2, y2, w2, h2, name2 = shapes_list[j]

                # Check if shapes overlap
                if (x1 < x2 + w2 and x1 + w1 > x2 and
                    y1 < y2 + h2 and y1 + h1 > y2):
                    overlapping_shapes.append((name1, name2))

        # Report findings
        if shapes_issues:
            print("Issues found:")
            for issue in shapes_issues:
                print(f"  [ERROR] {issue}")
        else:
            print("[OK] No boundary, font size, or table sizing issues found")

        if overlapping_shapes:
            print("\nOverlapping shapes found:")
            for name1, name2 in overlapping_shapes:
                print(f"  [WARNING]  {name1} overlaps with {name2}")
        else:
            print("\n[OK] No overlapping shapes found")

        print(f"\nTotal shapes on slide: {len([s for s in slide.shapes])}")
        print(f"Text/Table shapes: {len([s for s in slide.shapes if s.has_text_frame or s.has_table])}")

if __name__ == "__main__":
    pptx_file = r"C:/Users/DaudDewan/OneDrive - SymphonyAI/Documents/Learning/roehampton/deep_learning/la-reconvla/LA_ReconVLA_Presentation.pptx"
    analyze_presentation(pptx_file)